"""
exporters.py - Convert analysis results to PDF and PPTX formats.

Uses:
  - WeasyPrint: HTML to PDF with professional styling and page breaks
  - python-pptx: JSON to PPTX presentations with clean layouts
"""
import json
import os
from datetime import datetime
from pathlib import Path
from typing import Dict, List, Any

try:
    from weasyprint import HTML, CSS
except ImportError:
    HTML = CSS = None

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.dml.color import RGBColor
except ImportError:
    Presentation = None


# ============================================================================
# PDF EXPORTER (WeasyPrint)
# ============================================================================

def _generate_pdf_html(result: dict, domain: str) -> str:
    """Generate professional HTML for PDF export with page breaks."""
    g = result.get("graph_stats", {})
    anchors = result.get("anchors", {})
    clusters = result.get("clusters", {}).get("clusters", [])
    link_cands = result.get("link_candidates", [])

    # Executive summary metrics
    metrics = [
        ("Total Pages", g.get("pages_total", 0)),
        ("Indexable Pages", g.get("pages_indexable", 0)),
        ("Internal Links", g.get("internal_links", 0)),
        ("Avg Inlinks per Page", g.get("avg_inlinks", 0)),
        ("Max Crawl Depth", g.get("max_crawl_depth", 0)),
        ("Topical Clusters", len(clusters)),
    ]

    issues = [
        ("Orphan Pages", g.get("orphan_pages", []), "high"),
        ("Broken Links", g.get("broken_internal_links", []), "high"),
        ("Over-Optimized Anchors", anchors.get("over_optimized_anchors", []), "medium"),
        ("Generic Anchors", anchors.get("generic_anchors", []), "low"),
    ]

    # Build HTML with CSS for PDF
    html = f"""<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="utf-8">
    <meta name="viewport" content="width=device-width, initial-scale=1">
    <title>Internal Linking Intelligence - {domain}</title>
    <style>
        * {{
            margin: 0;
            padding: 0;
            box-sizing: border-box;
        }}

        @page {{
            size: A4;
            margin: 2.5cm;
            @bottom-center {{
                content: "Page " counter(page) " of " counter(pages);
                font-size: 11px;
                color: #888;
            }}
        }}

        body {{
            font-family: 'Segoe UI', Tahoma, Geneva, Verdana, sans-serif;
            color: #1a1a1a;
            line-height: 1.6;
            background: white;
        }}

        .container {{
            max-width: 100%;
        }}

        h1 {{
            font-size: 32px;
            color: #1a1a1a;
            margin-top: 20px;
            margin-bottom: 8px;
            border-bottom: 3px solid #3b82f6;
            padding-bottom: 12px;
        }}

        h2 {{
            font-size: 22px;
            color: #1f2937;
            margin-top: 25px;
            margin-bottom: 15px;
            page-break-after: avoid;
            border-left: 4px solid #3b82f6;
            padding-left: 12px;
        }}

        h3 {{
            font-size: 16px;
            color: #374151;
            margin-top: 18px;
            margin-bottom: 12px;
            page-break-after: avoid;
        }}

        .header {{
            text-align: center;
            margin-bottom: 30px;
            page-break-after: avoid;
        }}

        .domain {{
            font-size: 18px;
            color: #6b7280;
            margin-bottom: 20px;
        }}

        .metrics {{
            display: grid;
            grid-template-columns: 1fr 1fr 1fr;
            gap: 15px;
            margin: 20px 0;
            page-break-inside: avoid;
        }}

        .metric-card {{
            background: #f3f4f6;
            border: 1px solid #e5e7eb;
            border-radius: 8px;
            padding: 15px;
            text-align: center;
        }}

        .metric-value {{
            font-size: 24px;
            font-weight: bold;
            color: #3b82f6;
            margin-bottom: 5px;
        }}

        .metric-label {{
            font-size: 12px;
            color: #6b7280;
            text-transform: uppercase;
            letter-spacing: 0.5px;
        }}

        table {{
            width: 100%;
            border-collapse: collapse;
            margin: 15px 0;
            page-break-inside: auto;
        }}

        tr {{
            page-break-inside: avoid;
        }}

        th {{
            background: #3b82f6;
            color: white;
            padding: 12px;
            text-align: left;
            font-size: 12px;
            font-weight: 600;
            text-transform: uppercase;
            letter-spacing: 0.5px;
        }}

        td {{
            padding: 10px 12px;
            border-bottom: 1px solid #e5e7eb;
            font-size: 13px;
        }}

        tr:nth-child(even) {{
            background: #f9fafb;
        }}

        .mono {{
            font-family: 'Courier New', monospace;
            font-size: 11px;
            color: #1f2937;
            word-break: break-all;
        }}

        .badge {{
            display: inline-block;
            padding: 4px 10px;
            border-radius: 4px;
            font-size: 11px;
            font-weight: 600;
            text-transform: uppercase;
            letter-spacing: 0.3px;
        }}

        .badge.high {{
            background: #fee2e2;
            color: #991b1b;
        }}

        .badge.medium {{
            background: #fef3c7;
            color: #92400e;
        }}

        .badge.low {{
            background: #dbeafe;
            color: #1e40af;
        }}

        .badge.success {{
            background: #dcfce7;
            color: #166534;
        }}

        .section {{
            margin-top: 30px;
            margin-bottom: 20px;
        }}

        .page-break {{
            page-break-after: always;
        }}

        .issue-list {{
            margin: 15px 0;
        }}

        .issue-item {{
            margin: 10px 0;
            padding: 10px;
            background: #f9fafb;
            border-left: 3px solid #ef4444;
            page-break-inside: avoid;
        }}

        .issue-count {{
            font-weight: bold;
            color: #ef4444;
        }}

        .summary-box {{
            background: #eff6ff;
            border: 1px solid #bfdbfe;
            border-radius: 6px;
            padding: 15px;
            margin: 15px 0;
            page-break-inside: avoid;
        }}

        .generated {{
            margin-top: 40px;
            text-align: center;
            font-size: 11px;
            color: #9ca3af;
            border-top: 1px solid #e5e7eb;
            padding-top: 15px;
        }}

        .note {{
            font-style: italic;
            color: #6b7280;
            font-size: 12px;
            margin: 10px 0;
        }}
    </style>
</head>
<body>
<div class="container">
    <!-- TITLE PAGE -->
    <div class="header">
        <h1>Internal Linking Intelligence Report</h1>
        <div class="domain">{domain}</div>
        <div class="note">Generated on {datetime.now().strftime('%B %d, %Y at %I:%M %p')}</div>
    </div>

    <!-- EXECUTIVE SUMMARY -->
    <div class="section">
        <h2>Executive Summary</h2>
        <div class="metrics">
"""

    for label, value in metrics:
        html += f"""
            <div class="metric-card">
                <div class="metric-value">{value:,}</div>
                <div class="metric-label">{label}</div>
            </div>
"""

    html += """
        </div>
    </div>

    <!-- ISSUES OVERVIEW -->
    <div class="section page-break">
        <h2>Issues & Opportunities</h2>
"""

    for issue_label, issues_list, severity in issues:
        count = len(issues_list)
        if count > 0:
            html += f"""
        <div class="issue-item">
            <strong>{issue_label}:</strong> <span class="issue-count">{count}</span>
            <span class="badge {severity}">{severity}</span>
        </div>
"""

    html += """
    </div>

    <!-- TOPICAL CLUSTERS -->
    <div class="section page-break">
        <h2>Topical Clusters & Authority</h2>
        <div class="summary-box">
            <strong>Overview:</strong> Pages are organized into topical clusters to identify authority hubs and
            distribution of internal link equity. "Hub" indicates a clear authority page; "Scattered" means
            distribution is more uniform.
        </div>
        <table>
            <thead>
                <tr>
                    <th>Cluster</th>
                    <th>Pages</th>
                    <th>Authority</th>
                    <th>Hub Page</th>
                    <th>Keywords</th>
                </tr>
            </thead>
            <tbody>
"""

    for cluster in clusters[:15]:  # First 15 clusters
        hub_url = cluster.get("hub_page", "N/A")
        authority = cluster.get("authority", "scattered")
        keywords = ", ".join(cluster.get("keywords", [])[:5])

        html += f"""
                <tr>
                    <td><strong>{cluster.get('name', 'Unknown')}</strong></td>
                    <td>{cluster.get('size', 0)}</td>
                    <td><span class="badge {'success' if authority == 'hub' else 'low'}">{authority}</span></td>
                    <td class="mono">{hub_url.split('/')[-1]}</td>
                    <td><small>{keywords}</small></td>
                </tr>
"""

    html += """
            </tbody>
        </table>
    </div>

    <!-- ANCHOR TEXT ANALYSIS -->
    <div class="section page-break">
        <h2>Anchor Text Analysis</h2>
        <div class="summary-box">
            <strong>Key Metrics:</strong>
"""

    html += f"""
            <br>• <strong>Generic Anchors:</strong> {len(anchors.get('generic_anchors', []))} 
            <br>• <strong>Empty/Image Anchors:</strong> {len(anchors.get('empty_or_image_only', []))}
            <br>• <strong>Over-Optimized Anchors:</strong> {len(anchors.get('over_optimized_anchors', []))}
            <br>• <strong>Total Internal Links:</strong> {anchors.get('total_internal_anchors', 0)}
        </div>

        <h3>Top Over-Optimized Anchors</h3>
        <table>
            <thead>
                <tr>
                    <th>Destination</th>
                    <th>Anchor Text</th>
                    <th>Count</th>
                    <th>% of Links</th>
                </tr>
            </thead>
            <tbody>
"""

    for item in anchors.get("over_optimized_anchors", [])[:10]:
        html += f"""
                <tr>
                    <td class="mono">{item['destination'].split('/')[-1]}</td>
                    <td><strong>{item['anchor']}</strong></td>
                    <td>{item['count']}</td>
                    <td>{int(item['share'] * 100)}%</td>
                </tr>
"""

    html += """
            </tbody>
        </table>
    </div>

    <!-- LINK RECOMMENDATIONS -->
    <div class="section page-break">
        <h2>Contextual Link Recommendations</h2>
        <div class="summary-box">
            <strong>Strategy:</strong> These recommendations identify topically-related pages that should be
            internally linked to strengthen topical authority and improve crawlability.
        </div>

        <table>
            <thead>
                <tr>
                    <th>From Page</th>
                    <th>Target Page</th>
                    <th>Suggested Anchor</th>
                    <th>Relevance</th>
                </tr>
            </thead>
            <tbody>
"""

    for rec in link_cands[:20]:  # First 20 recommendations
        source = rec.get("source", "").split("/")[-1]
        for cand in rec.get("candidates", [])[:2]:  # Max 2 candidates per source
            target = cand.get("target", "").split("/")[-1]
            anchor = cand.get("suggested_anchor", "")
            rel = cand.get("relatedness", 0)

            html += f"""
                <tr>
                    <td class="mono">{source}</td>
                    <td class="mono">{target}</td>
                    <td><strong>{anchor}</strong></td>
                    <td>{rel:.3f}</td>
                </tr>
"""

    html += """
            </tbody>
        </table>
    </div>

    <!-- FOOTER -->
    <div class="generated">
        <strong>Link Intel Suite</strong> | Internal Linking Analysis & Optimization Report<br>
        Standard library based deterministic analysis | No external API calls
    </div>
</div>
</body>
</html>
"""
    return html


def export_to_pdf(result: dict, domain: str, output_path: str) -> bool:
    """Export analysis results to PDF using WeasyPrint.
    
    Args:
        result: Analysis result dict from analyzer.analyze()
        domain: Website domain for title
        output_path: Path to write PDF file
    
    Returns:
        True if successful, False otherwise
    """
    if HTML is None or CSS is None:
        print("❌ WeasyPrint not installed. Run: pip install weasyprint")
        return False

    try:
        html_content = _generate_pdf_html(result, domain)
        HTML(string=html_content).write_pdf(output_path)
        print(f"✅ PDF exported to {output_path}")
        return True
    except Exception as e:
        print(f"❌ PDF export failed: {e}")
        return False


# ============================================================================
# PPTX EXPORTER (python-pptx)
# ============================================================================

def _add_title_slide(prs: Presentation, domain: str, result: dict) -> None:
    """Add title slide to presentation."""
    g = result.get("graph_stats", {})
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(59, 130, 246)  # Blue

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_p = title_frame.paragraphs[0]
    title_p.text = "Internal Linking Intelligence"
    title_p.font.size = Pt(54)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(255, 255, 255)

    # Domain
    domain_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.6), Inches(9), Inches(0.5))
    domain_frame = domain_box.text_frame
    domain_p = domain_frame.paragraphs[0]
    domain_p.text = domain
    domain_p.font.size = Pt(28)
    domain_p.font.color.rgb = RGBColor(219, 234, 254)

    # Metrics
    metrics_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(2.5))
    metrics_frame = metrics_box.text_frame
    metrics_frame.word_wrap = True
    
    metrics_text = f"""
{g.get('pages_indexable', 0):,} Indexable Pages  •  {g.get('internal_links', 0):,} Internal Links
{len(result.get('clusters', {}).get('clusters', []))} Topical Clusters  •  {len(result.get('link_candidates', []))} Link Recommendations
"""
    
    metrics_p = metrics_frame.paragraphs[0]
    metrics_p.text = metrics_text.strip()
    metrics_p.font.size = Pt(16)
    metrics_p.font.color.rgb = RGBColor(255, 255, 255)
    metrics_p.alignment = PP_ALIGN.CENTER

    # Date
    date_box = slide.shapes.add_textbox(Inches(0.5), Inches(7), Inches(9), Inches(0.4))
    date_frame = date_box.text_frame
    date_p = date_frame.paragraphs[0]
    date_p.text = datetime.now().strftime("%B %d, %Y")
    date_p.font.size = Pt(12)
    date_p.font.color.rgb = RGBColor(219, 234, 254)
    date_p.alignment = PP_ALIGN.CENTER


def _add_section_slide(prs: Presentation, title: str) -> Any:
    """Add a section title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(31, 41, 55)  # Dark gray

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_p = title_frame.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(44)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(59, 130, 246)

    return slide


def _add_content_slide(prs: Presentation, title: str) -> Any:
    """Add a content slide with standard layout."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)

    # Title bar
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = RGBColor(59, 130, 246)
    title_shape.line.color.rgb = RGBColor(59, 130, 246)

    title_frame = title_shape.text_frame
    title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    title_p = title_frame.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(32)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(255, 255, 255)
    title_p.space_before = Pt(6)
    title_p.space_after = Pt(6)

    return slide


def export_to_pptx(result: dict, domain: str, output_path: str) -> bool:
    """Export analysis results to PPTX using python-pptx.
    
    Args:
        result: Analysis result dict from analyzer.analyze()
        domain: Website domain for presentation
        output_path: Path to write PPTX file
    
    Returns:
        True if successful, False otherwise
    """
    if Presentation is None:
        print("❌ python-pptx not installed. Run: pip install python-pptx")
        return False

    try:
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)

        g = result.get("graph_stats", {})
        anchors = result.get("anchors", {})
        clusters = result.get("clusters", {}).get("clusters", [])
        link_cands = result.get("link_candidates", [])

        # Slide 1: Title
        _add_title_slide(prs, domain, result)

        # Slide 2: Executive Summary
        _add_section_slide(prs, "Executive Summary")

        slide = _add_content_slide(prs, "Key Metrics")
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.8))
        text_frame = content_box.text_frame
        text_frame.word_wrap = True

        metrics_text = f"""• Total Pages: {g.get('pages_total', 0):,}
• Indexable Pages: {g.get('pages_indexable', 0):,}
• Internal Links: {g.get('internal_links', 0):,}
• Average Inlinks/Page: {g.get('avg_inlinks', 0)}
• Max Crawl Depth: {g.get('max_crawl_depth', 0)}
• Topical Clusters: {len(clusters)}
• Link Recommendations: {len(link_cands)}
"""

        for line in metrics_text.strip().split("\n"):
            p = text_frame.add_paragraph()
            p.text = line
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(31, 41, 55)
            p.space_before = Pt(6)
            p.space_after = Pt(6)
            p.level = 0

        # Slide 3: Issues Overview
        slide = _add_content_slide(prs, "Issues & Opportunities")
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.8))
        text_frame = content_box.text_frame
        text_frame.word_wrap = True

        issues_data = [
            ("Orphan Pages", len(g.get("orphan_pages", [])), "high"),
            ("Broken Links", len(g.get("broken_internal_links", [])), "high"),
            ("Redirect Links", len(g.get("redirect_internal_links", [])), "medium"),
            ("Nofollow Links", len(g.get("nofollow_internal_links", [])), "medium"),
            ("Over-Optimized Anchors", len(anchors.get("over_optimized_anchors", [])), "medium"),
            ("Generic Anchors", len(anchors.get("generic_anchors", [])), "low"),
        ]

        for label, count, severity in issues_data:
            p = text_frame.add_paragraph()
            p.text = f"• {label}: {count}"
            p.font.size = Pt(16)
            p.font.color.rgb = RGBColor(239, 68, 68) if severity == "high" else RGBColor(31, 41, 55)
            p.space_before = Pt(4)
            p.space_after = Pt(4)

        # Slide 4: Topical Clusters
        _add_section_slide(prs, "Topical Clusters & Authority")

        slide = _add_content_slide(prs, "Cluster Overview")
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.8))
        text_frame = content_box.text_frame
        text_frame.word_wrap = True

        for cluster in clusters[:8]:
            p = text_frame.add_paragraph()
            p.text = f"{cluster['name']} ({cluster['size']} pages)"
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = RGBColor(59, 130, 246)
            p.space_before = Pt(4)
            p.space_after = Pt(2)

            p2 = text_frame.add_paragraph()
            authority = "Hub" if cluster.get("authority") == "hub" else "Scattered"
            p2.text = f"  Authority: {authority} | Hub: {cluster.get('hub_page', 'N/A').split('/')[-1]}"
            p2.font.size = Pt(11)
            p2.font.color.rgb = RGBColor(107, 114, 128)
            p2.space_after = Pt(8)
            p2.level = 1

        # Slide 5: Anchor Text Analysis
        _add_section_slide(prs, "Anchor Text Analysis")

        slide = _add_content_slide(prs, "Anchor Issues")
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.8))
        text_frame = content_box.text_frame
        text_frame.word_wrap = True

        anchor_summary = f"""Total Internal Anchors: {anchors.get('total_internal_anchors', 0):,}

Generic/Non-Descriptive: {len(anchors.get('generic_anchors', []))}
Empty or Image Only: {len(anchors.get('empty_or_image_only', []))}
Over-Optimized: {len(anchors.get('over_optimized_anchors', []))}

Recommendation: Use descriptive, keyword-rich anchor text that accurately 
reflects the target page's topic. Avoid over-optimization of any single anchor.
"""

        for line in anchor_summary.strip().split("\n"):
            p = text_frame.add_paragraph()
            if line.startswith("•") or line.startswith("Total") or ":" in line:
                p.text = line
                p.font.size = Pt(14)
                p.font.bold = "Recommendation" in line
            else:
                p.text = line
                p.font.size = Pt(12)
            p.font.color.rgb = RGBColor(31, 41, 55)
            p.space_before = Pt(2)
            p.space_after = Pt(2)

        # Slide 6: Link Recommendations
        _add_section_slide(prs, "Link Recommendations")

        slide = _add_content_slide(prs, "Top Opportunities")
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.8))
        text_frame = content_box.text_frame
        text_frame.word_wrap = True

        for i, rec in enumerate(link_cands[:5]):
            source = rec.get("source", "").split("/")[-1] or "homepage"
            candidates = rec.get("candidates", [])[:1]
            
            p = text_frame.add_paragraph()
            p.text = f"From: {source}"
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = RGBColor(59, 130, 246)
            p.space_before = Pt(4)
            p.space_after = Pt(2)

            for cand in candidates:
                target = cand.get("target", "").split("/")[-1]
                anchor = cand.get("suggested_anchor", "")
                rel = cand.get("relatedness", 0)

                p2 = text_frame.add_paragraph()
                p2.text = f"  → {target}"
                p2.font.size = Pt(11)
                p2.font.color.rgb = RGBColor(31, 41, 55)
                p2.space_after = Pt(2)
                p2.level = 1

                p3 = text_frame.add_paragraph()
                p3.text = f"     Anchor: \"{anchor}\" (relevance: {rel:.3f})"
                p3.font.size = Pt(10)
                p3.font.italic = True
                p3.font.color.rgb = RGBColor(107, 114, 128)
                p3.space_after = Pt(6)
                p3.level = 2

        # Slide 7: Recommendations
        slide = _add_content_slide(prs, "Strategic Recommendations")
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.8))
        text_frame = content_box.text_frame
        text_frame.word_wrap = True

        recommendations = [
            "Strengthen Hub Pages: Internal link authority pages heavily",
            "Fix Broken Links: Redirects or remove all 4xx/5xx internal links",
            "Diversify Anchors: Use varied, descriptive anchor text",
            "Connect Related Topics: Use suggested link recommendations",
            "Optimize Crawl Depth: Reduce orphan pages via internal linking",
            "Monitor Authority: Track hub page inlink growth quarterly",
        ]

        for i, rec in enumerate(recommendations, 1):
            p = text_frame.add_paragraph()
            p.text = f"{i}. {rec}"
            p.font.size = Pt(13)
            p.font.color.rgb = RGBColor(31, 41, 55)
            p.space_before = Pt(4)
            p.space_after = Pt(4)

        prs.save(output_path)
        print(f"✅ PPTX exported to {output_path}")
        return True

    except Exception as e:
        print(f"❌ PPTX export failed: {e}")
        return False


# ============================================================================
# PUBLIC API
# ============================================================================

def export_results(result: dict, output_dir: str, domain: str = "website",
                   formats: List[str] = None) -> Dict[str, bool]:
    """Export analysis results to multiple formats.
    
    Args:
        result: Analysis result dict from analyzer.analyze()
        output_dir: Directory to write export files
        domain: Domain name for titles
        formats: List of formats to export ("pdf", "pptx", or both)
    
    Returns:
        Dict with format name as key, success boolean as value
    
    Example:
        >>> result = analyzer.analyze("sample-export")
        >>> exports = exporters.export_results(result, "outputs", "example.com")
        >>> if exports["pdf"] and exports["pptx"]:
        ...     print("All exports complete!")
    """
    if formats is None:
        formats = ["pdf", "pptx"]

    os.makedirs(output_dir, exist_ok=True)
    results = {}

    for fmt in formats:
        if fmt.lower() == "pdf":
            pdf_path = os.path.join(output_dir, f"report_{domain.replace('.', '_')}.pdf")
            results["pdf"] = export_to_pdf(result, domain, pdf_path)
        elif fmt.lower() == "pptx":
            pptx_path = os.path.join(output_dir, f"report_{domain.replace('.', '_')}.pptx")
            results["pptx"] = export_to_pptx(result, domain, pptx_path)

    return results


if __name__ == "__main__":
    # Example usage
    import sys
    from linkintel import analyzer

    if len(sys.argv) < 2:
        print("Usage: python exporters.py <export_dir> [domain]")
        sys.exit(1)

    export_dir = sys.argv[1]
    domain = sys.argv[2] if len(sys.argv) > 2 else "website.com"

    # Run analysis
    print(f"🔍 Analyzing {export_dir}...")
    result = analyzer.analyze(export_dir)

    # Export results
    print(f"📊 Exporting results...")
    exports = export_results(result, "outputs", domain, ["pdf", "pptx"])

    if all(exports.values()):
        print("✅ All exports completed successfully!")
    else:
        print("⚠️  Some exports failed. Check dependencies:")
        print("   pip install weasyprint python-pptx")
